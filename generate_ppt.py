"""生成 Super Tutor 毕业答辩 PPTX 文件 — 瑞士现代风格（Swiss Modern），11 页幻灯片。

【功能说明】
使用 python-pptx 库生成学术答辩演示文稿，包含：
- 封面（标题 + 副标题）
- 项目背景与目标
- 系统架构图（文字版）
- 核心技术亮点（知识点 DAG + 诊断规则 + 苏格拉底追问）
- 数据库设计（6 表 ER 关系）
- 前端界面截图说明
- 技术栈与代码量统计
- 创新点总结
- 演示效果展示
- 未来展望
- 致谢页

样式：Swiss Modern 风格 — 深蓝色主色调（#1A3A8A）+ 白色背景 + Arial 字体

输出文件：super-tutor-defense.pptx（11 页）

【耦合关系】
- 不依赖 super_tutor/ 下的任何模块（独立工具脚本）
- 仅依赖 python-pptx（外部库）
- 被手动运行：python generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Colors ──────────────────────────────────────────────────
C_BG        = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_ALT    = RGBColor(0xF7, 0xF8, 0xFB)
C_TEXT      = RGBColor(0x11, 0x11, 0x11)
C_TEXT_MID  = RGBColor(0x55, 0x55, 0x55)
C_TEXT_LOW  = RGBColor(0x99, 0x99, 0x99)
C_ACCENT    = RGBColor(0x1A, 0x3A, 0x8A)  # deep navy
C_ACCENT_LT = RGBColor(0xE8, 0xEB, 0xF4)
C_RED       = RGBColor(0xCC, 0x22, 0x00)
C_RED_LT    = RGBColor(0xFA, 0xE8, 0xE5)
C_GREEN     = RGBColor(0x0D, 0x7C, 0x4B)
C_GREEN_LT  = RGBColor(0xE5, 0xF4, 0xED)
C_ORANGE    = RGBColor(0xC2, 0x70, 0x0A)
C_ORANGE_LT = RGBColor(0xFD, 0xF2, 0xE0)
C_BORDER    = RGBColor(0xE0, 0xE2, 0xE8)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

FONT_TITLE = "Arial"
FONT_BODY  = "Arial"
FONT_MONO  = "Consolas"

# ── Helpers ─────────────────────────────────────────────────
def set_slide_bg(slide, color):
    """Set the background fill color of a slide.

    Args:
        slide: pptx Slide object.
        color: RGBColor to use as the solid background fill.
    """
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text="", font_size=14,
                color=C_TEXT, bold=False, font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
                line_spacing=1.2):
    """Add a single-paragraph text box to a slide.

    Args:
        slide: pptx Slide object.
        left, top, width, height: Position and size in inches.
        text: The text content.
        font_size: Font size in points.
        color: RGBColor for the text.
        bold: Whether the text is bold.
        font_name: Font family name (e.g. "Arial").
        alignment: Text horizontal alignment (PP_ALIGN.LEFT/CENTER/RIGHT).
        line_spacing: Line height multiplier (1.0 = single, 1.2 = 120%).

    Returns:
        The pptx TextFrame object for further manipulation.
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def add_rich_textbox(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple styled runs within a single paragraph.

    Args:
        slide: pptx Slide object.
        left, top, width, height: Position and size in inches.
        runs: List of (text, font_size, color, bold, font_name) tuples.
            Each tuple creates one run with independent styling.
        alignment: Text horizontal alignment.

    Returns:
        The pptx TextFrame object.
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, (text, fs, clr, bld, fn) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.name = fn
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=C_TEXT, bold_first=False, font_name=FONT_BODY,
                  line_spacing=1.4, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs (one per line).

    Args:
        slide: pptx Slide object.
        left, top, width, height: Position and size in inches.
        lines: List of strings, each becomes a separate paragraph.
        font_size: Font size in points.
        color: RGBColor for text.
        bold_first: If True, the first paragraph is rendered bold.
        font_name: Font family name.
        line_spacing: Line height multiplier.
        alignment: Text horizontal alignment.

    Returns:
        The pptx TextFrame object.
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = (bold_first and i == 0)
        p.alignment = alignment
        p.space_after = Pt(2)
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def add_rect(slide, left, top, width, height, fill_color=None,
             border_color=None, border_width=Pt(0)):
    """Add a rectangle shape (used for cards, panels, backgrounds).

    Args:
        slide: pptx Slide object.
        left, top, width, height: Position and size in inches.
        fill_color: Optional RGBColor for solid fill. None = transparent.
        border_color: Optional RGBColor for border. None = no border.
        border_width: Border thickness in Pt (default Pt(0)).

    Returns:
        The pptx Shape object.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
        shape.line.fill.solid()
    return shape

def add_line(slide, left, top, width, height, color=C_BORDER, width_pt=1):
    """Add a horizontal or vertical separator line (thin filled rectangle).

    Args:
        slide: pptx Slide object.
        left, top, width, height: Position and size in inches.
            Use very small height for horizontal lines, small width for vertical.
        color: RGBColor for the line fill.
        width_pt: Unused (kept for API compatibility).

    Returns:
        The pptx Shape object.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_bar(slide):
    """Add the left-edge navy accent bar (visual branding element).

    Renders a thin vertical bar (0.08" × 7.5") at the left edge of every slide.
    """
    add_rect(slide, 0, 0, 0.08, 7.5, fill_color=C_ACCENT)

def add_page_num(slide, current, total):
    """Add a page number badge to the bottom-right corner of a slide.

    Renders as "01 / 11" in monospace grey text.

    Args:
        slide: pptx Slide object.
        current: Current page number (1-based).
        total: Total number of slides.
    """
    add_textbox(slide, 11.55, 7.0, 1.5, 0.4,
                f"{current:02d} / {total:02d}", font_size=10,
                color=C_TEXT_LOW, font_name=FONT_MONO, alignment=PP_ALIGN.RIGHT)

def add_slide_header(slide, num_label, title):
    """Render a standard slide header: section tag + title + bottom border line.

    Layout (y positions in inches):
        num_label at y=0.55 (small navy caps, e.g. "背景")
        title     at y=0.82 (16pt bold)
        separator at y=1.20 (thin grey horizontal rule)

    Args:
        slide: pptx Slide object.
        num_label: Section label text (e.g. "核心流程", "架构").
        title: Slide title text (e.g. "七步学习闭环").
    """
    add_textbox(slide, 0.85, 0.55, 2.5, 0.35, num_label,
                font_size=9, color=C_ACCENT, bold=True, font_name=FONT_TITLE)
    add_textbox(slide, 0.85, 0.82, 8, 0.4, title,
                font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)
    add_line(slide, 0.85, 1.2, 11.6, 0.008, color=C_BORDER)

def add_tag(slide, left, top, text, bg_color, text_color):
    """Render a small colored tag/badge pill (1.3" × 0.28").

    Used throughout slides to label sections (e.g. "LLM解析", "Kahn算法").

    Args:
        slide: pptx Slide object.
        left, top: Position in inches.
        text: Label text (7pt bold, centered).
        bg_color: Background fill color.
        text_color: Text and border color.

    Returns:
        The pptx Shape object.
    """
    shape = add_rect(slide, left, top, 1.3, 0.28, fill_color=bg_color, border_color=text_color, border_width=Pt(0.5))
    shape.text_frame.word_wrap = False
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER
    return shape

# ── Create Presentation ─────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 11

# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_line(slide, 1.25, 4.95, 2.2, 0.024, color=C_RED, width_pt=3)

add_textbox(slide, 1.25, 2.0, 10, 0.3, "毕业设计答辩 · 计算机科学与技术",
            font_size=10, color=C_ACCENT, bold=True, font_name=FONT_TITLE)
add_textbox(slide, 1.25, 2.5, 11, 1.4, "Super Tutor",
            font_size=56, color=C_ACCENT, bold=True, font_name=FONT_TITLE)
add_textbox(slide, 1.25, 3.5, 11, 0.7, "基于 LLM 的智能教学系统",
            font_size=56, color=C_TEXT, bold=True, font_name=FONT_TITLE)
add_textbox(slide, 1.25, 4.2, 10, 0.5, "从教材上传到苏格拉底追问的个性化学习完整闭环",
            font_size=18, color=C_TEXT_MID, font_name=FONT_BODY)

meta_y = 5.4
for i, (label, value) in enumerate([
    ("答辩人", "李某某"), ("指导教师", "王教授"), ("日期", "2026 年 6 月")
]):
    x = 1.25 + i * 3.2
    add_textbox(slide, x, meta_y, 3, 0.25, label, font_size=11, color=C_TEXT_LOW, font_name=FONT_BODY)
    add_textbox(slide, x, meta_y + 0.25, 3, 0.3, value, font_size=13, color=C_ACCENT, bold=True, font_name=FONT_BODY)

add_page_num(slide, 1, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: BACKGROUND & PAIN POINTS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "背景", "传统学习工具的困境")

problems = [
    ("01", "缺乏个性化诊断",
     "大部分学习工具采用「一刀切」模式，无法根据学生的实际掌握情况动态调整学习路径。学生不知道自己的薄弱环节在哪，只能盲目刷题。"),
    ("02", "知识依赖关系被忽略",
     "知识点之间存在前驱/后继关系（如微积分依赖函数基础），传统工具无法建模这种 DAG 依赖，导致学生「跳级学习」。"),
    ("03", "错题缺乏深度追问",
     "做错的题目只能看解析，无法追问「为什么」。学生被动接受答案，没有经历主动思考的认知过程。"),
    ("04", "学习计划靠人工排期",
     "制定复习计划需要人工评估每章的优先级和时长，缺乏量化的自动化排期算法。"),
]

for i, (num, title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = 0.85 + col * 6.0
    y = 1.5 + row * 2.2

    card = add_rect(slide, x, y, 5.6, 1.9, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
    add_textbox(slide, x + 0.3, y + 0.2, 1, 0.3, num, font_size=14, color=C_RED, bold=True, font_name=FONT_MONO)
    add_textbox(slide, x + 0.3, y + 0.55, 5, 0.3, title, font_size=15, color=C_TEXT, bold=True, font_name=FONT_TITLE)
    add_textbox(slide, x + 0.3, y + 0.95, 5, 0.8, desc, font_size=11, color=C_TEXT_MID, font_name=FONT_BODY)

# Solution strip
add_rect(slide, 0.85, 6.15, 11.6, 0.65, fill_color=C_ACCENT)
add_textbox(slide, 1.1, 6.25, 1.2, 0.35, "解决思路", font_size=8, color=RGBColor(0xAA, 0xB8, 0xD8), bold=True, font_name=FONT_TITLE)
add_textbox(slide, 2.4, 6.2, 9.8, 0.5,
            "利用大语言模型（LLM）自动提取知识点依赖图，结合拓扑排序、前置规则校准和苏格拉底追问，构建个性化学习闭环",
            font_size=12, color=C_WHITE, font_name=FONT_BODY)

add_page_num(slide, 2, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: 7-STEP WORKFLOW
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "核心流程", "七步学习闭环")

steps = [
    ("①", "📄", "上传教材", "PDF上传 / 文本粘贴"),
    ("②", "🧠", "AI提取知识点", "LLM解析 + 建立双向依赖"),
    ("③", "🔍", "诊断评估", "拓扑排序 + 3条前置规则"),
    ("④", "📋", "学习计划", "优先级公式自动排期"),
    ("⑤", "✏️", "多题型练习", "6种题型 · 混合批改"),
    ("⑥", "📖", "错题本", "按KP分组 · 可重做"),
    ("⑦", "💬", "苏格拉底追问", "L1→L2→L3 层级引导"),
]

for i, (num, icon, title, detail) in enumerate(steps):
    row = 0 if i < 4 else 1
    col = i if i < 4 else i - 4
    if row == 1:
        col += 1  # offset bottom row

    x = 0.6 + col * 2.35
    y = 1.6 + row * 2.6

    card = add_rect(slide, x, y, 2.15, 1.8, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
    add_textbox(slide, x + 0.15, y + 0.15, 1, 0.3, num, font_size=16, color=C_ACCENT, bold=True, font_name=FONT_TITLE)
    add_textbox(slide, x + 0.15, y + 0.45, 1.8, 0.3, icon, font_size=18)
    add_textbox(slide, x + 0.15, y + 0.8, 1.8, 0.3, title, font_size=12, color=C_TEXT, bold=True, font_name=FONT_TITLE)
    add_textbox(slide, x + 0.15, y + 1.1, 1.8, 0.55, detail, font_size=9, color=C_TEXT_LOW, font_name=FONT_BODY)

# Arrows between step 4 and 5 (visual connection)
for i in range(3):
    x = 2.85 + i * 2.35
    add_textbox(slide, x, 2.2, 0.3, 0.3, "→", font_size=14, color=C_ACCENT, bold=True, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
for i in range(3):
    x = 3.85 + i * 2.35
    add_textbox(slide, x, 4.8, 0.3, 0.3, "→", font_size=14, color=C_ACCENT, bold=True, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)

add_page_num(slide, 3, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: TECHNICAL ARCHITECTURE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "架构", "五层技术架构")

layers = [
    ("🖥", "前端层", "Streamlit 1.35+", "单页应用，按钮驱动，20个 session_state 键管理UI状态", C_ACCENT, C_ACCENT_LT),
    ("⚙️", "业务引擎层", "5个无状态Engine", "KnowledgeEngine · AssessmentEngine · PlanEngine · QuizEngine · SocraticEngine", C_ACCENT, C_ACCENT_LT),
    ("💾", "数据层", "SQLite 6表 + aiosqlite", "WAL模式异步I/O，6张核心表（materials·knowledge_points·questions·quiz_attempts·wrong_questions·study_plans）", C_ACCENT, C_ACCENT_LT),
    ("🤖", "AI层", "DeepSeek API", "OpenAI SDK兼容接口，3次指数退避重试（1s→2s→4s），解析超时180s/批改超时120s", C_GREEN, C_GREEN_LT),
    ("📐", "基础设施层", "Pydantic 2.0 + PyPDF2", "8个Pydantic数据模型 + 2个枚举类型 + 5个Prompt模板 + 3级异常体系", C_ORANGE, C_ORANGE_LT),
]

for i, (icon, name, tag, desc, accent_c, bg_c) in enumerate(layers):
    y = 1.5 + i * 1.0
    add_rect(slide, 0.85, y, 11.6, 0.85, fill_color=bg_c, border_color=C_BORDER, border_width=Pt(0.5))
    # Left color strip
    add_line(slide, 0.85, y, 0.06, 0.85, color=accent_c, width_pt=4)
    add_textbox(slide, 1.2, y + 0.15, 0.5, 0.5, icon, font_size=16)
    add_textbox(slide, 1.7, y + 0.22, 0.9, 0.35, name, font_size=13, color=C_TEXT, bold=True, font_name=FONT_TITLE)
    add_tag(slide, 2.7, y + 0.25, tag, bg_c, accent_c)
    add_textbox(slide, 4.3, y + 0.22, 7.8, 0.4, desc, font_size=10, color=C_TEXT_MID, font_name=FONT_BODY)

add_page_num(slide, 4, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: INNOVATION 1 — KP DAG
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "创新点①", "知识点依赖图（DAG）自动构建")

# Left card
add_rect(slide, 0.85, 1.5, 5.6, 4.5, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_tag(slide, 1.15, 1.7, "LLM解析", C_ACCENT_LT, C_ACCENT)
add_textbox(slide, 1.15, 2.15, 5, 0.35, "从教材到知识点", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)
add_multiline(slide, 1.15, 2.6, 5, 2.0, [
    "— PyPDF2 提取 PDF 纯文本",
    "— LLM 将内容拆分为知识点列表",
    "— 每个KP含：title·summary·content·keywords",
    "— difficulty（5级）·prerequisite_indices",
    "— 自动生成UUID并批量写入 knowledge_points 表",
], font_size=12, color=C_TEXT_MID, font_name=FONT_BODY)

# Right card
add_rect(slide, 6.8, 1.5, 5.6, 4.5, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_tag(slide, 7.1, 1.7, "Kahn算法", C_ACCENT_LT, C_ACCENT)
add_textbox(slide, 7.1, 2.15, 5, 0.35, "拓扑排序 & 双向关系", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)
add_multiline(slide, 7.1, 2.6, 5, 1.5, [
    "— prerequisite_indices → 自动建立双向关系",
    "— prerequisite_ids ↔ successor_ids",
    "— Kahn算法拓扑排序，确保前驱KP在先",
    "— 难度分级：beginner→easy→medium→hard→expert",
], font_size=12, color=C_TEXT_MID, font_name=FONT_BODY)

# DAG visualization
dag_y = 4.6
nodes = [("KP1", C_RED, C_RED_LT), ("KP2", C_ACCENT, C_ACCENT_LT),
         ("KP3", C_ACCENT, C_ACCENT_LT), ("KP4", C_GREEN, C_GREEN_LT)]
for i, (label, clr, bg) in enumerate(nodes):
    x = 7.4 + i * 1.3
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(dag_y), Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = clr
    shape.line.width = Pt(1.5)
    shape.line.fill.solid()
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = clr
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER
    if i < 3:
        add_textbox(slide, x + 0.75, dag_y + 0.15, 0.4, 0.3, "→", font_size=14, color=C_TEXT_LOW, bold=True, font_name=FONT_MONO)

add_textbox(slide, 7.1, 5.5, 5, 0.3, "前驱 → 后继的拓扑序列", font_size=10, color=C_TEXT_LOW, font_name=FONT_BODY,
            alignment=PP_ALIGN.CENTER)

add_page_num(slide, 5, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 6: INNOVATION 2 — ASSESSMENT RULES
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "创新点②", "诊断评估与前置规则校准")

# Left: 3 rules
add_rect(slide, 0.85, 1.5, 5.6, 4.5, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_tag(slide, 1.15, 1.7, "3条前置规则", C_RED_LT, C_RED)
add_textbox(slide, 1.15, 2.15, 5, 0.35, "规则驱动的掌握度校准", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)

rules = [
    ("❶", "置信度折扣 — 前驱KP掌握度 ≤ 0.5 时，后继KP置信度 × 0.7，掌握度重新计算"),
    ("❷", "需复习标记 — 后继正确（≥ 0.6）但前驱错误（< 0.5）时，标记前驱为 need_review"),
    ("❸", "需重学标记 — ≥ 3个直接后继全部答错时，标记前驱为 need_relearn，掌握度折半"),
]
for i, (num, desc) in enumerate(rules):
    y = 2.7 + i * 1.0
    add_textbox(slide, 1.15, y, 0.4, 0.3, num, font_size=18, color=C_ACCENT, bold=True, font_name=FONT_TITLE)
    add_textbox(slide, 1.55, y + 0.05, 4.5, 0.7, desc, font_size=12, color=C_TEXT_MID, font_name=FONT_BODY)

# Right: mastery states
add_rect(slide, 6.8, 1.5, 5.6, 4.5, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_tag(slide, 7.1, 1.7, "掌握度分类", C_GREEN_LT, C_GREEN)
add_textbox(slide, 7.1, 2.15, 5, 0.35, "四级掌握度状态", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)

states = [
    ("mastered", "掌握度 ≥ 0.8", C_GREEN, C_GREEN_LT),
    ("learning", "掌握度 0.5 – 0.8", C_ACCENT, C_ACCENT_LT),
    ("need_review", "规则2触发或单一薄弱", C_ORANGE, C_ORANGE_LT),
    ("need_relearn", "≤ 0.3 或规则3触发", C_RED, C_RED_LT),
]
for i, (label, desc, clr, bg) in enumerate(states):
    y = 2.7 + i * 0.85
    add_rect(slide, 7.4, y, 4.7, 0.65, fill_color=bg, border_color=clr, border_width=Pt(0.5))
    add_textbox(slide, 7.6, y + 0.1, 1.8, 0.3, label, font_size=13, color=clr, bold=True, font_name=FONT_MONO)
    add_textbox(slide, 9.5, y + 0.12, 2.5, 0.3, desc, font_size=11, color=C_TEXT_MID, font_name=FONT_BODY)

add_page_num(slide, 6, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 7: INNOVATION 3 — GRADING + WRONG BOOK
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "创新点③", "混合批改与错题本")

# Left: grading table
add_rect(slide, 0.85, 1.5, 5.6, 4.5, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_tag(slide, 1.15, 1.7, "6种题型", C_ACCENT_LT, C_ACCENT)
add_textbox(slide, 1.15, 2.15, 5, 0.35, "程序化 + LLM 混合批改", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)

table_data = [
    ("题型", "前端控件", "批改方式", True),
    ("multiple_choice", "st.radio", "程序化比对", False),
    ("true_false", "st.radio", "程序化比对", False),
    ("fill_in_blank", "st.text_input", "LLM语义匹配", False),
    ("short_answer", "st.text_area", "LLM采分点评分", False),
    ("essay", "st.text_area", "LLM多维度评分", False),
    ("coding", "st.text_area", "LLM测试用例判定", False),
]
table = slide.shapes.add_table(7, 3, Inches(1.15), Inches(2.7), Inches(4.8), Inches(2.8)).table
for r, (c1, c2, c3, is_header) in enumerate(table_data):
    for c, text in enumerate([c1, c2, c3]):
        cell = table.cell(r, c)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10 if not is_header else 10)
        p.font.color.rgb = C_WHITE if is_header else C_TEXT_MID
        p.font.bold = is_header
        p.font.name = FONT_MONO if not is_header else FONT_TITLE
        p.alignment = PP_ALIGN.CENTER
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_ACCENT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_BG if r % 2 == 0 else C_BG_ALT

# Right: wrong book
add_rect(slide, 6.8, 1.5, 5.6, 4.5, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_tag(slide, 7.1, 1.7, "错题本", C_RED_LT, C_RED)
add_textbox(slide, 7.1, 2.15, 5, 0.35, "智能错题管理", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)
add_multiline(slide, 7.1, 2.7, 5, 2.5, [
    "— 按KP分组展示（st.expander折叠面板）",
    "— 去重逻辑：同一题多次错→递增attempt_count",
    "— 每道错题：题干·错误答案·正确答案·解析·犯错次数",
    "— 筛选：按知识点 / 按时间（全部·7天·30天）",
    "— 一键重做：为该KP生成3道新题",
    "— 苏格拉底追问入口：🗨 按钮直达对话",
], font_size=12, color=C_TEXT_MID, font_name=FONT_BODY)

add_page_num(slide, 7, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 8: INNOVATION 4 — SOCRATIC
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "创新点④", "苏格拉底式追问")

# Left: state machine
add_rect(slide, 0.85, 1.5, 5.6, 4.5, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_tag(slide, 1.15, 1.7, "层级状态机", C_ACCENT_LT, C_ACCENT)
add_textbox(slide, 1.15, 2.15, 5, 0.35, "Prompt驱动 + Python硬切", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)

state_colors = [
    ("L1_GUIDING  笼统引导", C_ACCENT, C_ACCENT_LT),
    ("L2_HINTING  具体提示", C_ORANGE, C_ORANGE_LT),
    ("L3_NEAR_ANSWER  接近答案", C_RED, C_RED_LT),
    ("RESOLVED  已解决 ✓", C_GREEN, C_GREEN_LT),
]
for i, (label, clr, bg) in enumerate(state_colors):
    x = 1.4 + i * 2.7
    y = 3.0
    shape = add_rect(slide, x, y, 2.4, 0.55, fill_color=bg, border_color=clr, border_width=Pt(1))
    p = shape.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = clr
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER
    if i < 3:
        add_textbox(slide, x + 2.45, y + 0.1, 0.25, 0.3, "→", font_size=14, color=C_TEXT_LOW, bold=True, font_name=FONT_MONO)

# Down arrows to SHOW_ANSWER
for i in range(4):
    add_textbox(slide, 2.2 + i * 2.7, 3.65, 0.5, 0.3, "↓", font_size=12, color=C_TEXT_LOW, bold=True, font_name=FONT_MONO,
                alignment=PP_ALIGN.CENTER)

shape = add_rect(slide, 3.5, 4.1, 5.5, 0.5, fill_color=RGBColor(0xF0, 0xF0, 0xF3), border_color=C_BORDER, border_width=Pt(0.5))
p = shape.text_frame.paragraphs[0]
p.text = "SHOW_ANSWER  显示答案 — 任一阶段均可触发"
p.font.size = Pt(9)
p.font.color.rgb = C_TEXT_MID
p.font.bold = True
p.font.name = FONT_TITLE
p.alignment = PP_ALIGN.CENTER

# Right: safety protections
add_rect(slide, 6.8, 1.5, 5.6, 4.5, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_tag(slide, 7.1, 1.7, "安全保护", C_RED_LT, C_RED)
add_textbox(slide, 7.1, 2.15, 5, 0.35, "两个程序化硬切", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)

# Trigger 1
add_rect(slide, 7.4, 2.8, 4.7, 1.3, fill_color=C_RED_LT, border_color=C_RED, border_width=Pt(0.5))
add_textbox(slide, 7.6, 2.9, 4.3, 0.3, "触发条件 1：关键词检测", font_size=13, color=C_RED, bold=True, font_name=FONT_TITLE)
add_textbox(slide, 7.6, 3.3, 4.3, 0.6, '学生输入含"显示答案/告诉我答案/我不会/太难了/想不出来"等关键词 → 直接SHOW_ANSWER',
            font_size=11, color=C_TEXT_MID, font_name=FONT_BODY)

# Trigger 2
add_rect(slide, 7.4, 4.3, 4.7, 1.3, fill_color=C_ORANGE_LT, border_color=C_ORANGE, border_width=Pt(0.5))
add_textbox(slide, 7.6, 4.4, 4.3, 0.3, "触发条件 2：最大轮数限制", font_size=13, color=C_ORANGE, bold=True, font_name=FONT_TITLE)
add_textbox(slide, 7.6, 4.8, 4.3, 0.6, "对话轮数 ≥ 6轮仍无法 RESOLVED → 强制 SHOW_ANSWER，避免无限追问",
            font_size=11, color=C_TEXT_MID, font_name=FONT_BODY)

add_page_num(slide, 8, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 9: PRIORITY FORMULA
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "算法", "学习计划优先级")

# Left: formula
add_rect(slide, 0.85, 1.5, 5.6, 2.0, fill_color=C_ACCENT_LT, border_color=C_ACCENT, border_width=Pt(1.5))
add_textbox(slide, 1.15, 1.7, 5, 0.25, "优先级公式", font_size=9, color=C_ACCENT, bold=True, font_name=FONT_TITLE)
add_textbox(slide, 1.15, 2.1, 5, 1.0,
            "priority = (1 - mastery)\n         × (1 + successor_count / total_kps)",
            font_size=18, color=C_TEXT, bold=False, font_name=FONT_MONO)

# Schedule box
add_rect(slide, 0.85, 3.8, 5.6, 2.4, fill_color=C_RED_LT, border_color=C_RED, border_width=Pt(1.5))
add_textbox(slide, 1.15, 4.0, 5, 0.25, "排期算法", font_size=9, color=C_RED, bold=True, font_name=FONT_TITLE)
add_multiline(slide, 1.15, 4.4, 5, 1.5, [
    "① Kahn拓扑排序 → KP序列",
    "② 优先级排序（相同拓扑层内）",
    "③ daily_schedule: start_date + n天",
    "④ 每日排期一个KP",
    "⑤ 估算时长: 难度×掌握度缺口（10-120分钟/KP）",
], font_size=12, color=C_TEXT_MID, font_name=FONT_BODY)

# Right: activity types
add_textbox(slide, 7.1, 1.5, 5, 0.35, "活动类型分配", font_size=16, color=C_TEXT, bold=True, font_name=FONT_TITLE)

activities = [
    ("< 0.3", "learn_new", "新学 — 从零开始学习", C_RED),
    ("0.3–0.5", "review", "复习 — 回顾薄弱环节", C_ORANGE),
    ("0.5–0.8", "practice", "练习 — 针对性强化训练", C_ACCENT),
    ("≥ 0.8", "quiz", "测验 — 巩固已掌握内容", C_GREEN),
]
for i, (range_, name, desc, clr) in enumerate(activities):
    y = 2.1 + i * 0.9
    add_rect(slide, 7.1, y, 5.4, 0.7, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
    add_textbox(slide, 7.3, y + 0.12, 0.8, 0.3, range_, font_size=12, color=clr, bold=True, font_name=FONT_MONO)
    add_textbox(slide, 8.15, y + 0.12, 1.2, 0.3, name, font_size=12, color=clr, bold=True, font_name=FONT_TITLE)
    add_textbox(slide, 9.4, y + 0.14, 2.8, 0.3, desc, font_size=11, color=C_TEXT_MID, font_name=FONT_BODY)

# Intent note
add_textbox(slide, 7.1, 6.2, 5.4, 0.8,
            "设计意图：掌握度越低 + 后继KP越多 → 优先级越高。\nsuccessor_count/total_kps 确保「被依赖越多」的知识点优先学习。",
            font_size=10, color=C_TEXT_LOW, font_name=FONT_BODY)

add_page_num(slide, 9, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 10: SUMMARY
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG)
add_accent_bar(slide)
add_slide_header(slide, "总结", "项目成果与展望")

stats = [
    ("~7,555", "行 Python 代码"),
    ("145", "个测试用例（全部通过）"),
    ("5+1", "个Engine + Streamlit前端"),
    ("6", "张SQLite数据表"),
]

for i, (num, label) in enumerate(stats):
    x = 0.85 + i * 3.0
    add_rect(slide, x, 1.6, 2.7, 1.8, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
    add_textbox(slide, x + 0.2, 1.9, 2.3, 0.8, num, font_size=38, color=C_ACCENT, bold=True, font_name=FONT_TITLE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 2.8, 2.3, 0.4, label, font_size=11, color=C_TEXT_MID, font_name=FONT_BODY,
                alignment=PP_ALIGN.CENTER)

# Future box
add_rect(slide, 0.85, 3.8, 11.6, 1.8, fill_color=C_BG_ALT, border_color=C_BORDER, border_width=Pt(0.5))
add_textbox(slide, 1.15, 4.0, 2, 0.3, "未来方向", font_size=14, color=C_ACCENT, bold=True, font_name=FONT_TITLE)
add_textbox(slide, 1.15, 4.5, 11, 0.8,
            "多学生支持（用户系统）· 学习分析仪表盘（数据可视化）· 更多题型/语言支持 · 语音交互苏格拉底追问 · 学习进度持久化与跨设备同步 · 向量检索增强知识点解析",
            font_size=12, color=C_TEXT_MID, font_name=FONT_BODY)

add_page_num(slide, 10, TOTAL)

# ══════════════════════════════════════════════════════════════
# SLIDE 11: CLOSING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_ACCENT)
add_rect(slide, 0, 0, 0.08, 7.5, fill_color=RGBColor(0x44, 0x55, 0xAA))

add_textbox(slide, 0, 2.6, 13.333, 1.2, "感谢聆听",
            font_size=60, color=C_WHITE, bold=True, font_name=FONT_TITLE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 3.9, 13.333, 0.6, "敬请各位老师批评指正",
            font_size=20, color=RGBColor(0xAA, 0xB8, 0xD8), font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 11.55, 7.0, 1.5, 0.4, "11 / 11", font_size=10,
            color=RGBColor(0x88, 0x99, 0xBB), font_name=FONT_MONO, alignment=PP_ALIGN.RIGHT)

# ── Save ────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "super-tutor-defense.pptx")
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
